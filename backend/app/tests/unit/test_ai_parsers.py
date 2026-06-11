"""多格式文档解析器纯单元测试。不依赖 conftest / DB、可 --noconftest 运行。"""

from __future__ import annotations

from io import BytesIO

import pytest
from docx import Document
from openpyxl import Workbook
from pptx import Presentation
from pptx.util import Inches

from app.core.document_state import DocumentStateError, DocumentStateMachine
from app.modules.ai.exceptions import DocumentParseError
from app.modules.ai.parsers import (
    MAX_EXTRACTED_TEXT_LENGTH,
    MAX_XLSX_ROWS,
    extract_text_from_bytes,
)


def _build_pdf_bytes(text: str) -> bytes:
    """手工构造含单页可提取文本的最小合法 PDF。"""
    content_stream = f"BT /F1 12 Tf 72 720 Td ({text}) Tj ET".encode("ascii")
    objects = [
        b"1 0 obj\n<< /Type /Catalog /Pages 2 0 R >>\nendobj\n",
        b"2 0 obj\n<< /Type /Pages /Kids [3 0 R] /Count 1 >>\nendobj\n",
        (
            b"3 0 obj\n<< /Type /Page /Parent 2 0 R /MediaBox [0 0 612 792] "
            b"/Resources << /Font << /F1 4 0 R >> >> /Contents 5 0 R >>\nendobj\n"
        ),
        b"4 0 obj\n<< /Type /Font /Subtype /Type1 /BaseFont /Helvetica >>\nendobj\n",
        (
            b"5 0 obj\n<< /Length "
            + str(len(content_stream)).encode("ascii")
            + b" >>\nstream\n"
            + content_stream
            + b"\nendstream\nendobj\n"
        ),
    ]
    output = bytearray(b"%PDF-1.4\n")
    offsets: list[int] = []
    for obj in objects:
        offsets.append(len(output))
        output += obj
    xref_position = len(output)
    output += b"xref\n0 6\n0000000000 65535 f \n"
    for offset in offsets:
        output += f"{offset:010d} 00000 n \n".encode("ascii")
    output += (
        b"trailer\n<< /Size 6 /Root 1 0 R >>\nstartxref\n"
        + str(xref_position).encode("ascii")
        + b"\n%%EOF\n"
    )
    return bytes(output)


def _build_docx_bytes(paragraph: str, table_cell: str) -> bytes:
    document = Document()
    document.add_paragraph(paragraph)
    table = document.add_table(rows=1, cols=1)
    table.cell(0, 0).text = table_cell
    buffer = BytesIO()
    document.save(buffer)
    return buffer.getvalue()


def _build_xlsx_bytes(sheet_title: str, rows: list[list[str]]) -> bytes:
    workbook = Workbook()
    worksheet = workbook.active
    assert worksheet is not None
    worksheet.title = sheet_title
    for row in rows:
        worksheet.append(row)
    buffer = BytesIO()
    workbook.save(buffer)
    return buffer.getvalue()


def _build_pptx_bytes(slide_text: str) -> bytes:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    textbox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(6), Inches(1))
    textbox.text_frame.text = slide_text
    buffer = BytesIO()
    presentation.save(buffer)
    return buffer.getvalue()


class TestPlainTextFormats:
    def test_txt_utf8_extracts_keyword(self) -> None:
        result = extract_text_from_bytes("员工手册正文内容".encode(), "txt")
        assert "员工手册正文内容" in result

    def test_txt_gb18030_extracts_keyword(self) -> None:
        result = extract_text_from_bytes("国标编码文档".encode("gb18030"), "txt")
        assert "国标编码文档" in result

    def test_md_extracts_keyword(self) -> None:
        result = extract_text_from_bytes(b"# Heading\nmarkdown body text", "md")
        assert "markdown body text" in result

    def test_csv_extracts_keyword(self) -> None:
        result = extract_text_from_bytes("列名,金额\n合同编号KU01,100".encode(), "csv")
        assert "合同编号KU01" in result


class TestBinaryFormats:
    def test_pdf_extracts_page_text(self) -> None:
        content = _build_pdf_bytes("Knowledge Uploader PDF Sample")
        result = extract_text_from_bytes(content, "pdf")
        assert "Knowledge Uploader PDF Sample" in result

    def test_docx_extracts_paragraph_and_table_text(self) -> None:
        content = _build_docx_bytes("段落关键词甲", "表格单元格乙")
        result = extract_text_from_bytes(content, "docx")
        assert "段落关键词甲" in result
        assert "表格单元格乙" in result

    def test_xlsx_extracts_sheet_name_and_cell_text(self) -> None:
        content = _build_xlsx_bytes("财务数据", [["合同编号", "金额"], ["KU-2026-001", "12345"]])
        result = extract_text_from_bytes(content, "xlsx")
        assert "财务数据" in result
        assert "KU-2026-001" in result

    def test_xlsx_rows_beyond_limit_are_skipped(self) -> None:
        rows = [[f"row-{index}"] for index in range(MAX_XLSX_ROWS + 5)]
        rows.append(["beyond-limit-marker"])
        content = _build_xlsx_bytes("大表", rows)
        result = extract_text_from_bytes(content, "xlsx")
        assert "row-0" in result
        assert "beyond-limit-marker" not in result

    def test_pptx_extracts_slide_text(self) -> None:
        content = _build_pptx_bytes("幻灯片关键词丙")
        result = extract_text_from_bytes(content, "pptx")
        assert "幻灯片关键词丙" in result

    def test_extension_lookup_is_case_insensitive(self) -> None:
        content = _build_pdf_bytes("Upper Case Extension")
        result = extract_text_from_bytes(content, "PDF")
        assert "Upper Case Extension" in result


class TestTruncationAndFallback:
    def test_long_text_truncated_to_max_length(self) -> None:
        content = ("a" * (MAX_EXTRACTED_TEXT_LENGTH + 5000)).encode("utf-8")
        result = extract_text_from_bytes(content, "txt")
        assert len(result) == MAX_EXTRACTED_TEXT_LENGTH

    def test_unknown_extension_returns_empty(self) -> None:
        assert extract_text_from_bytes(b"whatever", "xyz") == ""


class TestParseFailures:
    def test_corrupt_pdf_raises_parse_error_with_format(self) -> None:
        with pytest.raises(DocumentParseError) as exc_info:
            extract_text_from_bytes(b"not a real file", "pdf")
        assert "pdf" in str(exc_info.value)
        assert exc_info.value.format == "pdf"
        assert "not a real file" not in str(exc_info.value)

    def test_corrupt_docx_raises_parse_error_with_format(self) -> None:
        with pytest.raises(DocumentParseError) as exc_info:
            extract_text_from_bytes(b"not a real file", "docx")
        assert "docx" in str(exc_info.value)

    def test_corrupt_xlsx_raises_parse_error_with_format(self) -> None:
        with pytest.raises(DocumentParseError) as exc_info:
            extract_text_from_bytes(b"not a real file", "xlsx")
        assert "xlsx" in str(exc_info.value)

    def test_corrupt_pptx_raises_parse_error_with_format(self) -> None:
        with pytest.raises(DocumentParseError) as exc_info:
            extract_text_from_bytes(b"not a real file", "pptx")
        assert "pptx" in str(exc_info.value)

    @pytest.mark.parametrize(
        ("legacy", "upgrade"),
        [("doc", "docx"), ("xls", "xlsx"), ("ppt", "pptx")],
    )
    def test_legacy_format_raises_parse_error_with_upgrade_hint(
        self, legacy: str, upgrade: str
    ) -> None:
        with pytest.raises(DocumentParseError) as exc_info:
            extract_text_from_bytes(b"legacy binary", legacy)
        message = str(exc_info.value)
        assert legacy in message
        assert upgrade in message
        assert "重新上传" in message


class TestDocumentStateMachineRetry:
    def test_analysis_failed_can_retry_to_extracting_text(self) -> None:
        result = DocumentStateMachine.transition("analysis_failed", "extracting_text")
        assert result == "extracting_text"

    def test_analyzed_cannot_go_back_to_extracting_text(self) -> None:
        with pytest.raises(DocumentStateError):
            DocumentStateMachine.transition("analyzed", "extracting_text")


class TestDocumentLifecycleTransitions:
    @pytest.mark.parametrize(
        "from_status",
        [
            "uploaded",
            "pending_review",
            "approved",
            "rejected",
            "failed",
            "parsed",
            "analysis_failed",
            "analyzed",
            "sensitive_review_required",
            "disabled",
        ],
    )
    def test_soft_delete_allowed_from_stable_statuses(self, from_status: str) -> None:
        assert DocumentStateMachine.transition(from_status, "deleted") == "deleted"

    @pytest.mark.parametrize(
        "from_status",
        ["queued", "syncing", "uploaded_to_ragflow", "parsing", "extracting_text", "analyzing"],
    )
    def test_soft_delete_rejected_from_mid_pipeline_statuses(self, from_status: str) -> None:
        with pytest.raises(DocumentStateError):
            DocumentStateMachine.transition(from_status, "deleted")

    @pytest.mark.parametrize(
        "from_status",
        ["approved", "parsed", "failed", "rejected", "analyzed", "pending_review"],
    )
    def test_archive_allowed_from_stable_statuses(self, from_status: str) -> None:
        assert DocumentStateMachine.transition(from_status, "disabled") == "disabled"

    @pytest.mark.parametrize("from_status", ["uploaded", "syncing", "deleted"])
    def test_archive_rejected_from_other_statuses(self, from_status: str) -> None:
        with pytest.raises(DocumentStateError):
            DocumentStateMachine.transition(from_status, "disabled")

    def test_deleted_is_terminal(self) -> None:
        for to_status in ("uploaded", "pending_review", "approved", "disabled"):
            with pytest.raises(DocumentStateError):
                DocumentStateMachine.transition("deleted", to_status)

    def test_analyzed_can_reset_to_analysis_failed_for_reanalysis(self) -> None:
        assert DocumentStateMachine.transition("analyzed", "analysis_failed") == "analysis_failed"

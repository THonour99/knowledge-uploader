"""按扩展名分发的文档文本解析器注册表。

所有解析器签名统一为 ``Callable[[bytes], str]``。解析失败时抛 ``DocumentParseError``。
错误信息只含格式名与原因、绝不包含文件内容。
"""

from __future__ import annotations

import re
from collections.abc import Callable
from io import BytesIO
from zipfile import BadZipFile

from docx import Document as load_docx
from docx.opc.exceptions import PackageNotFoundError as DocxPackageNotFoundError
from openpyxl import load_workbook
from openpyxl.utils.exceptions import InvalidFileException
from pptx import Presentation as load_pptx
from pptx.exc import PackageNotFoundError as PptxPackageNotFoundError
from pypdf import PdfReader
from pypdf.errors import PyPdfError

from .exceptions import DocumentParseError

MAX_EXTRACTED_TEXT_LENGTH = 20_000
MAX_PDF_PAGES = 200
MAX_XLSX_ROWS = 2000
TEXT_ENCODINGS = ("utf-8", "utf-8-sig", "gb18030")
LEGACY_FORMAT_UPGRADES = {"doc": "docx", "xls": "xlsx", "ppt": "pptx"}

Parser = Callable[[bytes], str]
ExtractedTable = dict[str, object]


def parse_plain_text(content: bytes) -> str:
    for encoding in TEXT_ENCODINGS:
        try:
            return content.decode(encoding).strip()
        except UnicodeDecodeError:
            continue
    return content.decode("utf-8", errors="ignore").strip()


def parse_pdf(
    content: bytes,
    *,
    max_pages: int = MAX_PDF_PAGES,
    max_chars: int = MAX_EXTRACTED_TEXT_LENGTH,
) -> str:
    try:
        reader = PdfReader(BytesIO(content), strict=False)
        parts: list[str] = []
        total_length = 0
        for index, page in enumerate(reader.pages):
            if index >= max_pages or total_length >= max_chars:
                break
            text = (page.extract_text() or "").strip()
            if text:
                parts.append(text)
                total_length += len(text)
        return "\n".join(parts)
    except (PyPdfError, ValueError, KeyError, OSError) as exc:
        raise DocumentParseError(
            format="pdf",
            reason=f"文件损坏或内容无法读取 ({type(exc).__name__})",
        ) from exc


def parse_docx(content: bytes) -> str:
    try:
        document = load_docx(BytesIO(content))
        parts: list[str] = [
            paragraph.text.strip() for paragraph in document.paragraphs if paragraph.text.strip()
        ]
        for table in document.tables:
            for row in table.rows:
                cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                if cells:
                    parts.append("\t".join(cells))
        return "\n".join(parts)
    except (DocxPackageNotFoundError, BadZipFile, KeyError, ValueError) as exc:
        raise DocumentParseError(
            format="docx",
            reason=f"文件损坏或内容无法读取 ({type(exc).__name__})",
        ) from exc


def parse_xlsx(content: bytes) -> str:
    try:
        workbook = load_workbook(BytesIO(content), read_only=True, data_only=True)
        try:
            parts: list[str] = []
            total_rows = 0
            for worksheet in workbook.worksheets:
                parts.append(f"[{worksheet.title}]")
                for row in worksheet.iter_rows(values_only=True):
                    if total_rows >= MAX_XLSX_ROWS:
                        break
                    cells = [
                        str(value).strip()
                        for value in row
                        if value is not None and str(value).strip()
                    ]
                    if cells:
                        parts.append("\t".join(cells))
                    total_rows += 1
                if total_rows >= MAX_XLSX_ROWS:
                    break
            return "\n".join(parts)
        finally:
            workbook.close()
    except (BadZipFile, InvalidFileException, KeyError, ValueError, OSError) as exc:
        raise DocumentParseError(
            format="xlsx",
            reason=f"文件损坏或内容无法读取 ({type(exc).__name__})",
        ) from exc


def extract_tables_from_bytes(
    content: bytes,
    extension: str,
    *,
    max_pages: int = MAX_PDF_PAGES,
) -> list[ExtractedTable]:
    """提取可离线识别的表格结构, 输出 JSON 友好的表头、行、列和 Markdown."""

    normalized = extension.lower().lstrip(".")
    if normalized in LEGACY_FORMAT_UPGRADES:
        upgrade = LEGACY_FORMAT_UPGRADES[normalized]
        raise DocumentParseError(
            format=normalized,
            reason=f"不支持的旧格式。请转存为 {upgrade} 后重新上传",
        )
    if normalized == "xlsx":
        return extract_xlsx_tables(content)
    if normalized == "docx":
        return extract_docx_tables(content)
    if normalized == "pdf":
        return extract_pdf_tables(content, max_pages=max_pages)
    return []


def extract_xlsx_tables(content: bytes) -> list[ExtractedTable]:
    try:
        workbook = load_workbook(BytesIO(content), read_only=True, data_only=True)
        try:
            tables: list[ExtractedTable] = []
            total_rows = 0
            for worksheet in workbook.worksheets:
                rows: list[list[str]] = []
                for row in worksheet.iter_rows(values_only=True):
                    if total_rows >= MAX_XLSX_ROWS:
                        break
                    rows.append([_cell_to_text(value) for value in row])
                    total_rows += 1
                table = _build_table(title=worksheet.title, raw_rows=rows)
                if table is not None:
                    tables.append(table)
                if total_rows >= MAX_XLSX_ROWS:
                    break
            return tables
        finally:
            workbook.close()
    except (BadZipFile, InvalidFileException, KeyError, ValueError, OSError) as exc:
        raise DocumentParseError(
            format="xlsx",
            reason=f"文件损坏或内容无法读取 ({type(exc).__name__})",
        ) from exc


def extract_docx_tables(content: bytes) -> list[ExtractedTable]:
    try:
        document = load_docx(BytesIO(content))
        tables: list[ExtractedTable] = []
        for table in document.tables:
            rows = [[cell.text.strip() for cell in row.cells] for row in table.rows]
            table_data = _build_table(title=None, raw_rows=rows)
            if table_data is not None:
                tables.append(table_data)
        return tables
    except (DocxPackageNotFoundError, BadZipFile, KeyError, ValueError) as exc:
        raise DocumentParseError(
            format="docx",
            reason=f"文件损坏或内容无法读取 ({type(exc).__name__})",
        ) from exc


def extract_pdf_tables(
    content: bytes,
    *,
    max_pages: int = MAX_PDF_PAGES,
) -> list[ExtractedTable]:
    """基于 pypdf 文本层的离线表格兜底识别。

    当前依赖清单没有 pdfplumber, 因此只识别文本层中带管道、制表符、逗号或多空格
    分隔的连续行。扫描版或坐标型 PDF 表格留给 OCR/PDF 表格引擎扩展。
    """

    try:
        reader = PdfReader(BytesIO(content), strict=False)
        tables: list[ExtractedTable] = []
        for page_index, page in enumerate(reader.pages):
            if page_index >= max_pages:
                break
            text = page.extract_text() or ""
            for group_index, rows in enumerate(_table_rows_from_text(text), start=1):
                table = _build_table(
                    title=f"Page {page_index + 1} Table {group_index}",
                    raw_rows=rows,
                )
                if table is not None:
                    tables.append(table)
        return tables
    except (PyPdfError, ValueError, KeyError, OSError) as exc:
        raise DocumentParseError(
            format="pdf",
            reason=f"文件损坏或内容无法读取 ({type(exc).__name__})",
        ) from exc


def append_tables_markdown(
    text: str,
    tables: list[ExtractedTable],
    *,
    max_chars: int = MAX_EXTRACTED_TEXT_LENGTH,
) -> str:
    markdown_parts = [
        markdown
        for table in tables
        if isinstance((markdown := table.get("markdown")), str) and markdown.strip()
    ]
    if not markdown_parts:
        return text[:max_chars]
    combined = "\n\n".join(part for part in [text.strip(), *markdown_parts] if part)
    return combined[:max_chars]


def parse_pptx(content: bytes) -> str:
    try:
        presentation = load_pptx(BytesIO(content))
        parts: list[str] = []
        for slide in presentation.slides:
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                text = shape.text_frame.text.strip()
                if text:
                    parts.append(text)
        return "\n".join(parts)
    except (PptxPackageNotFoundError, BadZipFile, KeyError, ValueError) as exc:
        raise DocumentParseError(
            format="pptx",
            reason=f"文件损坏或内容无法读取 ({type(exc).__name__})",
        ) from exc


PARSER_REGISTRY: dict[str, Parser] = {
    "txt": parse_plain_text,
    "md": parse_plain_text,
    "csv": parse_plain_text,
    "pdf": parse_pdf,
    "docx": parse_docx,
    "xlsx": parse_xlsx,
    "pptx": parse_pptx,
}


def extract_text_from_bytes(
    content: bytes,
    extension: str,
    *,
    max_pages: int = MAX_PDF_PAGES,
    max_chars: int = MAX_EXTRACTED_TEXT_LENGTH,
) -> str:
    """按扩展名解析文本。

    截断上限默认沿用模块常量; 调用方 (如 ai.service 的异步任务) 可在
    调用前读取 runtime_config 的 processing.parse_max_pages /
    processing.parse_max_chars 并显式传入, 本函数保持纯同步无 IO。
    """
    normalized = extension.lower().lstrip(".")
    if normalized in LEGACY_FORMAT_UPGRADES:
        upgrade = LEGACY_FORMAT_UPGRADES[normalized]
        raise DocumentParseError(
            format=normalized,
            reason=f"不支持的旧格式。请转存为 {upgrade} 后重新上传",
        )
    parser = PARSER_REGISTRY.get(normalized)
    if parser is None:
        return ""
    if parser is parse_pdf:
        # parse_pdf 内部仅按页粗截断 (末页可溢出 max_chars), 外层切片是统一出口的精确硬上限
        return parse_pdf(content, max_pages=max_pages, max_chars=max_chars)[:max_chars]
    return parser(content)[:max_chars]


def _cell_to_text(value: object) -> str:
    if value is None:
        return ""
    return str(value).strip()


def _build_table(title: str | None, raw_rows: list[list[str]]) -> ExtractedTable | None:
    rows = _normalize_rows(raw_rows)
    if not rows:
        return None
    column_count = max(len(row) for row in rows)
    headers = _pad_row(rows[0], column_count)
    if not any(headers):
        headers = [f"Column {index}" for index in range(1, column_count + 1)]
    body_rows = [_pad_row(row, column_count) for row in rows[1:]]
    columns = [
        [row[column_index] for row in body_rows if row[column_index]]
        for column_index in range(column_count)
    ]
    markdown = _table_to_markdown(headers=headers, rows=body_rows)
    table: ExtractedTable = {
        "headers": headers,
        "rows": body_rows,
        "columns": columns,
        "markdown": markdown,
    }
    if title:
        table["title"] = title
    return table


def _normalize_rows(raw_rows: list[list[str]]) -> list[list[str]]:
    rows: list[list[str]] = []
    for raw_row in raw_rows:
        row = [_clean_cell(cell) for cell in raw_row]
        while row and not row[-1]:
            row.pop()
        if any(row):
            rows.append(row)
    return rows


def _clean_cell(value: str) -> str:
    return re.sub(r"\s+", " ", value).strip()


def _pad_row(row: list[str], column_count: int) -> list[str]:
    if len(row) >= column_count:
        return row
    return [*row, *([""] * (column_count - len(row)))]


def _table_to_markdown(*, headers: list[str], rows: list[list[str]]) -> str:
    header_line = "| " + " | ".join(_escape_markdown_cell(cell) for cell in headers) + " |"
    separator_line = "| " + " | ".join("---" for _ in headers) + " |"
    body_lines = [
        "| " + " | ".join(_escape_markdown_cell(cell) for cell in row) + " |" for row in rows
    ]
    return "\n".join([header_line, separator_line, *body_lines])


def _escape_markdown_cell(value: str) -> str:
    return value.replace("|", "\\|").replace("\n", " ").strip()


def _table_rows_from_text(text: str) -> list[list[list[str]]]:
    groups: list[list[list[str]]] = []
    current: list[list[str]] = []
    for line in text.splitlines():
        row = _split_table_line(line)
        if row is None:
            if len(current) >= 2:
                groups.append(current)
            current = []
            continue
        current.append(row)
    if len(current) >= 2:
        groups.append(current)
    return groups


def _split_table_line(line: str) -> list[str] | None:
    stripped = line.strip().strip("|").strip()
    if not stripped:
        return None
    if "|" in stripped:
        cells = stripped.split("|")
    elif "\t" in stripped:
        cells = stripped.split("\t")
    elif "," in stripped:
        cells = stripped.split(",")
    elif re.search(r"\s{2,}", stripped):
        cells = re.split(r"\s{2,}", stripped)
    else:
        return None
    row = [_clean_cell(cell) for cell in cells]
    return row if sum(1 for cell in row if cell) >= 2 else None
